# -*- coding: utf-8 -*-
"""Генерация презентации INMYHEART для защиты выпускного проекта."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT = PROJECT_ROOT / "Презентация_INMYHEART.pptx"

# Цвета: медицинская тема
COLOR_PRIMARY = RGBColor(0x0D, 0x47, 0xA1)  # тёмно-синий
COLOR_ACCENT = RGBColor(0x00, 0x96, 0x88)   # бирюза
COLOR_TEXT = RGBColor(0x21, 0x21, 0x21)
COLOR_MUTED = RGBColor(0x55, 0x55, 0x55)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(1.15),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(9), Inches(0.4))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = COLOR_MUTED


def _add_bullets(slide, items: list[str], *, top: float = 1.75, left: float = 0.6, width: float = 8.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)


def _add_numbered(slide, items: list[str], *, top: float = 1.75) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(1.2))
    tp = title.text_frame.paragraphs[0]
    tp.text = "AI-ассистент клиники INMYHEART"
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_PRIMARY

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(2.9), Inches(8.6), Inches(1.5))
    stf = sub.text_frame
    lines = [
        "RAG · OpenAI · ChromaDB · FastAPI · Telegram · VPS",
        "Ответы пациентам на организационные вопросы по базе знаний",
    ]
    for i, line in enumerate(lines):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if i == 0 else 16)
        p.font.color.rgb = COLOR_MUTED if i else COLOR_ACCENT

    foot = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(8), Inches(0.5))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Выпускной проект · 2026 · github.com/DimonRonD/InMyHeart"
    fp.font.size = Pt(12)
    fp.font.color.rgb = COLOR_MUTED


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_title_bar(slide, "Проблема")
    _add_bullets(
        slide,
        [
            "Сеть из 20 офисов: ~150 обращений на регистратуру за 12-часовую смену (1 запрос / 5 мин).",
            "До 3 000 обращений в день — в основном повторяющиеся справочные вопросы.",
            "Запись, подготовка к анализам, расписание, цены, правила посещения — одни и те же ответы.",
            "Администраторы перегружены: очередь на линии, задержки, риск несогласованных формулировок.",
            "Медицинские и персональные вопросы нельзя автоматизировать «в лоб» — нужны ограничения и эскалация.",
        ],
    )


def slide_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_title_bar(slide, "Решение", "Что делает ассистент")
    _add_bullets(
        slide,
        [
            "RAG-ассистент отвечает только по документам клиники (48 файлов, ~242 чанка в ChromaDB).",
            "Организационные вопросы (ORG) — краткий ответ + указание источника из базы знаний.",
            "ПДн (PII) и медицинские темы (MED) — отказ или эскалация на оператора в Telegram.",
            "Quality check: программная проверка + LLM post-check перед отправкой ответа.",
            "Каналы MVP: Telegram-бот, REST API (FastAPI /ask), CLI для тестирования.",
        ],
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_title_bar(slide, "Как работает система", "Стек: OpenAI · VPS · Telegram")
    _add_bullets(
        slide,
        [
            "1. Классификатор (gpt-5.4-nano) → ORG / PII / MED / RESULTS / OTHER",
            "2. ORG → семантический поиск Chroma + rerank → генерация ответа (gpt-5.4-mini)",
            "3. Проверка качества → ответ пациенту или перевод на оператора",
            "4. Forum Topics: отдельная ветка на клиента, relay «оператор ↔ пациент»",
            "5. SQLite: логи диалогов (без ПДн), KPI времени ответа, SLA ≤ 5 с",
            "",
            "Инфра: VPS (FastAPI + bot + Chroma + SQLite) · Telegram Bot API · OpenAI Embeddings",
        ],
        top=1.55,
    )


def slide_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_title_bar(slide, "Демо", "Сценарий для защиты (test_questions.md)")

    left_items = [
        "Типовые (ORG + RAG):",
        "• Как записаться к терапевту через Telegram?",
        "• До скольки работает клиника в среду?",
        "• Как подготовиться к крови на гормоны?",
        "• Сколько стоит УЗИ брюшной полости?",
        "• Когда принимает кардиолог Петров?",
        "",
        "Вне темы:",
        "• «Иванов Петр, ул. Ленина 15…» → PII, отказ",
        "• «Болит грудь, тяжело дышать» → MED, эскалация",
    ]
    _add_bullets(slide, left_items, top=1.55, left=0.55, width=5.2)

    demo_box = slide.shapes.add_shape(1, Inches(5.9), Inches(1.7), Inches(3.5), Inches(2.8))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    demo_box.line.color.rgb = COLOR_PRIMARY

    dtf = demo_box.text_frame
    dtf.word_wrap = True
    dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = dtf.paragraphs[0]
    p1.text = "Демо"
    p1.font.bold = True
    p1.font.size = Pt(20)
    p1.font.color.rgb = COLOR_PRIMARY
    p1.alignment = PP_ALIGN.CENTER

    for line in [
        "Telegram-бот",
        "или",
        "python scripts/assistant_cli.py",
        "",
        "← вставьте @username бота",
    ]:
        p = dtf.add_paragraph()
        p.text = line
        p.font.size = Pt(13 if "username" in line else 14)
        p.font.color.rgb = COLOR_MUTED if "username" in line else COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER


def slide_metrics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_title_bar(slide, "Результаты и метрики")

    metrics = [
        ("База знаний", "48 документов · ~242 чанка · txt, md, pdf, docx, csv"),
        ("Точность тестов", "pytest: 23/23 · AI-оценка CLI --test: ~87% PASS"),
        ("Время ответа", "SLA ≤ 5 с (требование курса MVP: < 30 с) ✓"),
        ("Маршрутизация", "ORG / PII / MED / RESULTS / OTHER + эскалация"),
        ("Автоматизация", "~70% обращений → 630 000 AI-диалогов / год"),
        ("CI", "GitHub Actions · pytest integration"),
    ]

    top = 1.65
    for label, value in metrics:
        lbl = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(2.4), Inches(0.45))
        lp = lbl.text_frame.paragraphs[0]
        lp.text = label
        lp.font.bold = True
        lp.font.size = Pt(16)
        lp.font.color.rgb = COLOR_ACCENT

        val = slide.shapes.add_textbox(Inches(3.1), Inches(top), Inches(6.3), Inches(0.55))
        vp = val.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(16)
        vp.font.color.rgb = COLOR_TEXT
        top += 0.72


def slide_costs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_title_bar(slide, "Экономика", "Расчёт на 12 месяцев · costs.md")

    _add_bullets(
        slide,
        [
            "Нагрузка: 20 офисов × 150 обращений/смену × 300 дней = 900 000 обращений/год.",
            "Реалистичная доля через AI: 70% → 630 000 LLM-диалогов/год.",
            "",
            "Рекомендуемый сценарий (гибрид gpt-5.4-nano + gpt-5.4-mini):",
            "  • OpenAI API — ~190 500 ₽/год",
            "  • VPS + backup — ~46 400 ₽/год",
            "  • Telegram, Chroma, SQLite — 0 ₽",
            "  • Резерв 12% — ~28 400 ₽",
            "",
            "ИТОГО: ~265 000 ₽/год (~22 000 ₽/мес) · ~0,30 ₽ за AI-диалог",
            "Экономия времени администраторов: ~31 500 человеко-часов/год",
        ],
        top=1.55,
    )


def slide_future(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_title_bar(slide, "Планы на будущее")
    _add_numbered(
        slide,
        [
            "Веб-чат для пациентов (UI поверх FastAPI /ask).",
            "Замена синтетической базы на актуальные регламенты клиники INMYHEART.",
            "Продакшен: Docker, Telegram webhook, мониторинг и алерты по KPI.",
            "Интеграция с CRM / МИС, авторизация операторов.",
            "Relay медиа (фото, файлы), расширение аналитики KPI.",
        ],
    )


def slide_thanks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_PRIMARY)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(8.6), Inches(1))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Спасибо за внимание!"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_WHITE
    tp.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(8.6), Inches(1.5))
    stf = sub.text_frame
    for i, line in enumerate(
        [
            "INMYHEART · AI-ассистент на RAG",
            "GitHub: github.com/DimonRonD/InMyHeart",
            "Готов ответить на вопросы и показать демо",
        ]
    ):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p.alignment = PP_ALIGN.CENTER


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_demo(prs)
    slide_metrics(prs)
    slide_costs(prs)
    slide_future(prs)
    slide_thanks(prs)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
